"""Exporta presentacion.html a presentacion-editable.pptx con FORMAS NATIVAS.

A diferencia de exportar_pptx.py (que rasteriza cada diapositiva con Chrome y la pega como
imagen, fiel al diseno pero no editable), este script reconstruye cada diapositiva con cuadros
de texto, vinetas y tablas de PowerPoint. Se puede abrir en Google Slides y editar. El diseno
queda aproximado: el HTML usa CSS Grid, pildoras redondeadas y unidades vh que PowerPoint no
tiene, asi que los chips y las barras se simplifican a texto con fondo.

    .venv/bin/python entrega/presentacion/exportar_pptx_editable.py [--fuentes mac|portable|office]

Los graficos van como imagen suelta (se pueden mover y redimensionar). Los SVG se rasterizan
con Chrome al vuelo. El guion.md sigue entrando como notas del orador.
"""
import argparse
import base64
import math
import re
import shutil
import subprocess
import tempfile
from html.parser import HTMLParser
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

AQUI = Path(__file__).resolve().parent
HTML = AQUI / 'presentacion.html'
GUION = AQUI / 'guion.md'
SALIDA = AQUI / 'presentacion-editable.pptx'

ANCHO, ALTO = Inches(13.333), Inches(7.5)
MARGEN, TOPE, PIE = Inches(0.62), Inches(0.42), Inches(0.42)
UTIL = ANCHO - 2 * MARGEN

C = {'ink': RGBColor(0x17, 0x24, 0x2C), 'ink2': RGBColor(0x3D, 0x4B, 0x5A),
     'muted': RGBColor(0x5F, 0x71, 0x83), 'line': RGBColor(0xD8, 0xDF, 0xE7),
     'acc': RGBColor(0x0E, 0x9B, 0x7E), 'acc2': RGBColor(0x70, 0x52, 0xC9),
     'mal': RGBColor(0xD4, 0x2A, 0x63), 'bg': RGBColor(0xF4, 0xF7, 0xF6),
     'card': RGBColor(0xFF, 0xFF, 0xFF), 'dark': RGBColor(0x17, 0x24, 0x2C),
     'claro': RGBColor(0xF4, 0xF7, 0xF6), 'verdeclaro': RGBColor(0xE4, 0xF2, 0xEE),
     'lila': RGBColor(0xED, 0xE9, 0xF8), 'gris': RGBColor(0xEF, 0xF3, 0xF5),
     'darkmuted': RGBColor(0x93, 0xA6, 0xB1), 'darkline': RGBColor(0x2E, 0x44, 0x50)}
# Un .pptx guarda UN nombre de fuente por run, no una pila de fallbacks como el CSS.
# 'portable': existe en macOS, Windows y Google Slides. 'mac': lo que ve el HTML en macOS.
FUENTES = {'portable': ('Arial', 'Courier New'),
           'mac': ('Helvetica', 'Menlo'),
           'office': ('Arial', 'Consolas')}
SANS, MONO = FUENTES['office']


# ------------------------------------------------------------------ parseo del HTML

class Nodo:
    def __init__(self, tag, attrs=None):
        self.tag, self.attrs, self.hijos, self.texto = tag, attrs or {}, [], ''

    def clases(self):
        return self.attrs.get('class', '').split()

    def tiene(self, c):
        return c in self.clases()

    def buscar(self, tag=None, clase=None):
        out = []
        for h in self.hijos:
            if isinstance(h, Nodo):
                if (tag is None or h.tag == tag) and (clase is None or h.tiene(clase)):
                    out.append(h)
                out.extend(h.buscar(tag, clase))
        return out

    def plano(self):
        """Texto sin marcas, con los espacios colapsados."""
        p = []
        for h in self.hijos:
            p.append(h if isinstance(h, str) else h.plano())
        return re.sub(r'\s+', ' ', ''.join(p)).strip()

    def runs(self, negrita=False, color=None):
        """[(texto, negrita, color, mono)] respetando <b>, <code>, <em>, .ok, .mal."""
        out = []
        for h in self.hijos:
            if isinstance(h, str):
                if h.strip():
                    out.append((re.sub(r'\s+', ' ', h), negrita, color, False))
                elif out:
                    out.append((' ', negrita, color, False))
                continue
            col = color
            if h.tiene('ok'):
                col = C['acc']
            elif h.tiene('mal'):
                col = C['mal']
            if h.tag == 'code':
                out.append((h.plano(), negrita, col or C['ink'], True))
                continue
            if h.tag == 'br':
                out.append(('\n', negrita, color, False))
                continue
            out.extend(h.runs(negrita or h.tag in ('b', 'strong'), col))
        return out


VACIOS = {'img', 'br', 'hr', 'meta', 'link', 'input', 'source'}


class Arbol(HTMLParser):
    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.raiz = Nodo('root')
        self.pila = [self.raiz]
        self.en_svg = 0
        self.svg = []

    def handle_starttag(self, tag, attrs):
        d = dict(attrs)
        if tag == 'svg' or self.en_svg:
            self.en_svg += 1
            self.svg.append(self.get_starttag_text())
            if self.en_svg == 1:
                n = Nodo('svgbox', d)
                self.pila[-1].hijos.append(n)
                self.pila.append(n)
            return
        if tag in VACIOS:
            self.pila[-1].hijos.append(Nodo(tag, d))
            return
        n = Nodo(tag, d)
        self.pila[-1].hijos.append(n)
        self.pila.append(n)

    def handle_startendtag(self, tag, attrs):
        if self.en_svg:
            self.svg.append(self.get_starttag_text())
            return
        self.pila[-1].hijos.append(Nodo(tag, dict(attrs)))

    def handle_endtag(self, tag):
        if self.en_svg:
            self.svg.append(f'</{tag}>')
            self.en_svg -= 1
            if self.en_svg == 0:
                self.pila[-1].texto = ''.join(self.svg)
                self.svg = []
                self.pila.pop()
            return
        if tag in VACIOS:
            return
        for i in range(len(self.pila) - 1, 0, -1):
            if self.pila[i].tag == tag:
                del self.pila[i:]
                return

    def handle_data(self, data):
        if self.en_svg:
            self.svg.append(data)
        elif data.strip() or data:
            self.pila[-1].hijos.append(data)


def leer_slides():
    html = HTML.read_text(encoding='utf-8')
    cuerpo = html[html.index('<section class="slide'):html.index('<div id="hud">')]
    a = Arbol()
    a.feed(cuerpo)
    return [n for n in a.raiz.hijos if isinstance(n, Nodo) and n.tag == 'section']


# ------------------------------------------------------------------ helpers de dibujo

def chrome():
    for b in ('google-chrome', 'google-chrome-stable', 'chromium', 'chromium-browser'):
        if shutil.which(b):
            return b
    for ruta in ('/Applications/Google Chrome.app/Contents/MacOS/Google Chrome',
                 '/Applications/Chromium.app/Contents/MacOS/Chromium'):
        if Path(ruta).exists():
            return ruta
    raise SystemExit('no encuentro google-chrome ni chromium')


def rasterizar_svg(svg, css, carpeta, nombre):
    """El SVG solo, sobre fondo claro, capturado a 2x. PowerPoint no lee SVG del HTML."""
    m = re.search(r'viewBox="0 0 ([\d.]+) ([\d.]+)"', svg)
    w, h = (float(m.group(1)), float(m.group(2))) if m else (1000.0, 700.0)
    esc = max(1, int(2200 / w))
    doc = carpeta / f'{nombre}.html'
    # el CSS de la presentacion le pone max-height en vh al svg: aca hay que anularlo
    doc.write_text(f'<meta charset="utf-8"><style>{css}\nhtml,body{{margin:0;'
                   f'background:#F4F7F6}}\n.figzoo svg,svg{{display:block !important;'
                   f'width:{w}px !important;height:{h}px !important;max-height:none !important;'
                   f'max-width:none !important}}</style>'
                   f'<div class="figzoo">{svg}</div>', encoding='utf-8')
    png = carpeta / f'{nombre}.png'
    subprocess.run([chrome(), '--headless=new', '--disable-gpu', '--hide-scrollbars',
                    f'--window-size={int(w)},{int(h)}', f'--force-device-scale-factor={esc}',
                    '--virtual-time-budget=4000', f'--screenshot={png}', f'file://{doc}'],
                   check=True, capture_output=True)
    return png, w / h


def texto(slide, x, y, w, h, runs, tam=12, color=None, alin=PP_ALIGN.LEFT,
          espacio=1.15, ancla=MSO_ANCHOR.TOP):
    caja = slide.shapes.add_textbox(x, y, w, h)
    tf = caja.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = ancla
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    parrafos = [[]]
    for t, neg, col, mono in runs:
        for i, trozo in enumerate(t.split('\n')):
            if i:
                parrafos.append([])
            if trozo:
                parrafos[-1].append((trozo, neg, col, mono))
    for i, ps in enumerate(parrafos):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alin
        p.line_spacing = espacio
        for t, neg, col, mono in ps:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(tam)
            r.font.bold = neg
            r.font.name = MONO if mono else SANS
            r.font.color.rgb = col or color or C['ink2']
    return caja


def sin_estilo(forma):
    """Saca el <p:style> del tema: si queda, PowerPoint y LibreOffice le meten sombra."""
    el = forma._element
    for st in el.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}style'):
        el.remove(st)
    forma.shadow.inherit = False
    return forma


def rect(slide, x, y, w, h, color, radio=False):
    """Rectangulo plano sin borde ni sombra: barras, reglas y acentos."""
    f = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radio else MSO_SHAPE.RECTANGLE, x, y, w, h)
    f.fill.solid()
    f.fill.fore_color.rgb = color
    f.line.fill.background()
    return sin_estilo(f)


def caja_texto(slide, x, y, w, h, relleno=None, borde=None, radio=True, izq=None):
    forma = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radio else MSO_SHAPE.RECTANGLE, x, y, w, h)
    forma.adjustments[0] = 0.06 if radio else 0
    if relleno is None:
        forma.fill.background()
    else:
        forma.fill.solid()
        forma.fill.fore_color.rgb = relleno
    if borde is None:
        forma.line.fill.background()
    else:
        forma.line.color.rgb = borde
        forma.line.width = Pt(1)
    sin_estilo(forma)
    forma.text_frame.text = ''
    if izq:  # barra de acento a la izquierda, como el border-left del HTML
        rect(slide, x, y, Inches(0.05), h, izq)
    return forma


def alto_texto(runs, ancho_in, tam, espacio=1.2):
    """Alto estimado: PowerPoint no reflowea al generar, hay que calcularlo."""
    chars = sum(len(t) for t, *_ in runs) or 1
    por_linea = max(10, int(ancho_in * 72 / (tam * 0.52)))
    lineas = sum(max(1, math.ceil(len(seg) / por_linea))
                 for seg in ''.join(t for t, *_ in runs).split('\n')) or 1
    lineas = max(lineas, math.ceil(chars / por_linea))
    return Inches(lineas * tam * espacio * 1.03 / 72)


# ------------------------------------------------------------------ bloques

def png_de_img(nodo, carpeta, nombre):
    src = nodo.attrs.get('src', '')
    m = re.match(r'data:image/(\w+);base64,(.*)$', src, re.S)
    if not m:
        return None, 1.0
    datos = base64.b64decode(m.group(2))
    ruta = carpeta / f'{nombre}.{m.group(1)}'
    ruta.write_bytes(datos)
    w, h = int.from_bytes(datos[16:20], 'big'), int.from_bytes(datos[20:24], 'big')
    return ruta, (w / h if h else 1.0)


def pintar(slide, nodo, x, y, w, ctx):
    """Dibuja un nodo del HTML y devuelve el alto consumido."""
    cls = nodo.clases()
    win = w / 914400

    if nodo.tag == 'svgbox':
        png, ratio = rasterizar_svg(nodo.texto, ctx['css'], ctx['tmp'], f'svg{ctx["n"]}')
        ctx['n'] += 1
        h = min(Inches(win / ratio), ctx['max_fig'])
        slide.shapes.add_picture(str(png), x + (w - Emu(int(h * ratio))) // 2, y,
                                 Emu(int(h * ratio)), h)
        return h

    if nodo.tag == 'img':
        png, ratio = png_de_img(nodo, ctx['tmp'], f'img{ctx["n"]}')
        ctx['n'] += 1
        if png is None:
            return Emu(0)
        h = min(Inches(win / ratio), ctx['max_fig'])
        slide.shapes.add_picture(str(png), x + (w - Emu(int(h * ratio))) // 2, y,
                                 Emu(int(h * ratio)), h)
        return h

    if nodo.tag == 'table':
        return pintar_tabla(slide, nodo, x, y, w)

    if nodo.tag == 'ul':
        alto = Emu(0)
        for li in [h for h in nodo.hijos if isinstance(h, Nodo) and h.tag == 'li']:
            runs = [('•  ', True, C['acc'], False)] + li.runs()
            hh = alto_texto(runs, win, 12) + Inches(0.09)
            texto(slide, x, y + alto, w, hh, runs, tam=12)
            alto += hh
        return alto

    if 'destacado' in cls:
        runs = nodo.runs()
        hh = alto_texto(runs, win - 0.5, 12.5) + Inches(0.28)
        caja_texto(slide, x, y, w, hh, C['verdeclaro'], None, True, C['acc'])
        texto(slide, x + Inches(0.2), y + Inches(0.13), w - Inches(0.35),
              hh - Inches(0.26), runs, tam=12.5, color=C['ink'])
        return hh + Inches(0.12)

    if 'heroe' in cls:  # numero grande arriba, glosa chica debajo (son display:block)
        b = nodo.buscar('b')
        sp = nodo.buscar('span')
        alto = Emu(0)
        if b:
            texto(slide, x, y, w, Inches(0.72), [(b[0].plano(), True, C['acc'], True)], tam=40,
                  espacio=1.0)
            alto += Inches(0.78)
        if sp:
            r = [(sp[0].plano(), False, C['muted'], False)]
            hh = alto_texto(r, win, 10.5) + Inches(0.06)
            texto(slide, x, y + alto, w, hh, r, tam=10.5)
            alto += hh
        return alto + Inches(0.2)

    if 'res' in cls:  # la flecha sale de un ::before del CSS, no esta en el DOM
        runs = [('→  ', True, C['acc'], False)] + nodo.runs()
        hh = alto_texto(runs, win, 12) + Inches(0.1)
        rect(slide, x, y, w, Inches(0.008), C['line'])
        texto(slide, x, y + Inches(0.12), w, hh, runs, tam=12, color=C['ink'])
        return hh + Inches(0.16)

    if 'kicker' in cls:
        em = nodo.buscar('em')
        etiqueta = nodo.plano()
        r = []
        if em:  # el <em> del kicker no va en versalitas ni en el color de acento
            etiqueta = etiqueta.replace(em[0].plano(), '').strip()
            r = [(etiqueta.upper() + '   ', True, C['acc2'], True),
                 (em[0].plano(), False, C['muted'], False)]
        else:
            r = [(etiqueta.upper(), True, C['acc2'], True)]
        texto(slide, x, y, w, Inches(0.2), r, tam=9)
        return Inches(0.3)

    if 'chips' in cls:
        piezas = []
        for ch in nodo.buscar('span'):
            if ch.tiene('chip'):
                piezas.append(ch)
        alto, cx, fila = Inches(0.0), x, Inches(0.30)
        for ch in piezas:
            t = ch.plano()
            cw = Inches(0.085 * len(t) + 0.22)
            if cx + cw > x + w:
                cx, alto = x, alto + fila + Inches(0.07)
            verde = ch.tiene('der')
            caja_texto(slide, cx, y + alto, cw, fila,
                       C['verdeclaro'] if verde else C['gris'],
                       C['acc'] if verde else C['line'])
            texto(slide, cx, y + alto, cw, fila, [(t, True, C['ink'], True)],
                  tam=9, alin=PP_ALIGN.CENTER, ancla=MSO_ANCHOR.MIDDLE)
            cx += cw + Inches(0.06)
        return alto + fila + Inches(0.12)

    if {'filaNums', 'flujo', 'grid4', 'tri', 'tres', 'linea', 'tira'} & set(cls):
        return pintar_grilla(slide, nodo, x, y, w, ctx)

    if 'panel' in cls or nodo.tag == 'div' and not cls:
        return pintar_contenedor(slide, nodo, x, y, w, ctx,
                                 marco='panel' in cls, alerta='alerta' in cls)

    if nodo.tag in ('p', 'h3', 'dl'):
        if nodo.tag == 'dl':
            return pintar_spec(slide, nodo, x, y, w)
        runs = nodo.runs()
        if not runs:
            return Emu(0)
        tam = 12.5 if 'grande' in cls else 12
        hh = alto_texto(runs, win, tam) + Inches(0.1)
        texto(slide, x, y, w, hh, runs, tam=tam)
        return hh

    return pintar_contenedor(slide, nodo, x, y, w, ctx)


def pintar_contenedor(slide, nodo, x, y, w, ctx, marco=False, alerta=False):
    pad = Inches(0.17) if marco else Inches(0)
    interior_x, interior_w = x + pad, w - 2 * pad
    alto = pad
    formas = len(slide.shapes._spTree)
    nodos = [h for h in nodo.hijos if isinstance(h, Nodo)]
    tope = ctx['max_fig']
    if len(nodos) > 1 and any(h.tag in ('img', 'svgbox') for h in nodos):
        ctx['max_fig'] = Emu(int(tope * 0.55))  # la figura comparte columna: dejar lugar abajo
    for h in nodo.hijos:
        if isinstance(h, Nodo):
            alto += pintar(slide, h, interior_x, y + alto, interior_w, ctx)
    ctx['max_fig'] = tope
    alto += pad
    if marco:
        caja = caja_texto(slide, x, y, w, alto, C['card'], C['line'], True,
                          C['mal'] if alerta else None)
        slide.shapes._spTree.remove(caja._element)
        slide.shapes._spTree.insert(formas, caja._element)  # el marco va detras del texto
    return alto + Inches(0.1)


def pintar_spec(slide, nodo, x, y, w):
    """<dl class="spec|specd"> -> dos columnas de termino y glosa."""
    oscuro = nodo.tiene('specd')
    pares, term = [], None
    for h in nodo.hijos:
        if isinstance(h, Nodo) and h.tag == 'dt':
            term = h.plano()
        elif isinstance(h, Nodo) and h.tag == 'dd' and term is not None:
            pares.append((term, h.plano()))
            term = None
    cols = 2 if oscuro and len(pares) > 6 else 1
    filas = math.ceil(len(pares) / cols)
    cw = w // cols
    fila = Inches(0.27)
    for i, (t, d) in enumerate(pares):
        cx = x + (i // filas) * cw
        cy = y + (i % filas) * fila
        texto(slide, cx, cy, Inches(1.7), fila, [(t, True, None, True)], tam=10,
              color=RGBColor(0xD6, 0xE4, 0xDF) if oscuro else C['muted'])
        texto(slide, cx + Inches(1.8), cy, cw - Inches(1.9), fila, [(d, False, None, False)],
              tam=10.5, color=RGBColor(0x81, 0x95, 0xA1) if oscuro else C['ink2'])
    return fila * filas + Inches(0.12)


SIN_ESTILO_TABLA = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'  # "sin estilo, sin cuadricula"
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def linea_inferior(celda, color, pt=1.0):
    """Borde de abajo: python-pptx no lo expone, hay que meterlo en el tcPr."""
    from pptx.oxml.ns import qn
    tc = celda._tc.get_or_add_tcPr()
    for viejo in tc.findall(qn('a:lnB')):
        tc.remove(viejo)
    ln = tc.makeelement(qn('a:lnB'), {'w': str(int(pt * 12700)), 'cap': 'flat',
                                      'cmpd': 'sng', 'algn': 'ctr'})
    fill = ln.makeelement(qn('a:solidFill'), {})
    c = ln.makeelement(qn('a:srgbClr'), {'val': f'{color}'})
    fill.append(c)
    ln.append(fill)
    tc.insert(0, ln)


def pintar_tabla(slide, nodo, x, y, w):
    filas = nodo.buscar('tr')
    if not filas:
        return Emu(0)
    ncol = max(len(f.buscar('th')) + len(f.buscar('td')) for f in filas)
    alto_fila = Inches(0.42)
    tabla = slide.shapes.add_table(len(filas), ncol, x, y, w, alto_fila * len(filas)).table
    from pptx.oxml.ns import qn
    tabla._tbl.tblPr.set('firstRow', '0')
    tabla._tbl.tblPr.set('bandRow', '0')
    for hijo in tabla._tbl.tblPr.findall(qn('a:tableStyleId')):
        tabla._tbl.tblPr.remove(hijo)
    est = tabla._tbl.tblPr.makeelement(qn('a:tableStyleId'), {})
    est.text = SIN_ESTILO_TABLA           # sin el tema azul de PowerPoint
    tabla._tbl.tblPr.append(est)

    for i, f in enumerate(filas):
        celdas = [c for c in f.hijos if isinstance(c, Nodo) and c.tag in ('td', 'th')]
        cab = bool(celdas) and celdas[0].tag == 'th'
        for j, c in enumerate(celdas[:ncol]):
            cel = tabla.cell(i, j)
            cel.margin_left = cel.margin_right = Inches(0.09)
            cel.margin_top = cel.margin_bottom = Inches(0.06)
            cel.vertical_anchor = MSO_ANCHOR.MIDDLE
            cel.fill.solid()
            cel.fill.fore_color.rgb = C['card']
            linea_inferior(cel, 'D8DFE7', 2.0 if cab else 1.0)
            propio = C['acc'] if c.tiene('ok') else (C['mal'] if c.tiene('mal') else None)
            p = cel.text_frame.paragraphs[0]
            for t, neg, col, mono in (c.runs() or [(c.plano(), False, None, False)]):
                r = p.add_run()
                r.text = t.upper() if cab else t
                r.font.size = Pt(9 if cab else 10.5)
                r.font.bold = neg or cab
                r.font.name = MONO if (mono or cab or c.tiene('n')) else SANS
                r.font.color.rgb = col or propio or (C['muted'] if cab else C['ink2'])
    return alto_fila * len(filas) + Inches(0.14)


def pintar_grilla(slide, nodo, x, y, w, ctx):
    hijos = [h for h in nodo.hijos if isinstance(h, Nodo)]
    hijos = [h for h in hijos if h.tag != 'div' or h.plano() or h.buscar('img')]
    if not hijos:
        return Emu(0)
    if nodo.tiene('flujo'):  # pasos anchos, flechas angostas
        flechas = [h for h in hijos if h.tiene('flecha')]
        pasos = [h for h in hijos if not h.tiene('flecha')]
        af = Inches(0.42)
        pw = (w - af * len(flechas) - Inches(0.1) * (len(hijos) - 1)) // max(1, len(pasos))
        cx, alto = x, Inches(1.35)
        for h in hijos:
            if h.tiene('flecha'):
                texto(slide, cx, y, af, alto, [('→', True, C['muted'], False)], tam=20,
                      alin=PP_ALIGN.CENTER, ancla=MSO_ANCHOR.MIDDLE)
                cx += af + Inches(0.1)
            else:
                acc = h.tiene('acc')
                caja_texto(slide, cx, y, pw, alto, C['verdeclaro'] if acc else C['card'],
                           C['acc'] if acc else C['line'])
                texto(slide, cx + Inches(0.1), y, pw - Inches(0.2), alto, h.runs(), tam=11.5,
                      alin=PP_ALIGN.CENTER, ancla=MSO_ANCHOR.MIDDLE, color=C['ink'])
                cx += pw + Inches(0.1)
        return alto + Inches(0.12)

    vertical = nodo.tiene('vert')
    cols = 1 if vertical else (2 if nodo.tiene('grid4') else len(hijos))
    alto_fila = Inches(1.55) if nodo.tiene('grid4') else Inches(1.1)
    filas = math.ceil(len(hijos) / cols)
    gap = Inches(0.16)
    cw = (w - gap * (cols - 1)) // cols
    altos = []
    for i, h in enumerate(hijos):
        cx = x + (i % cols) * (cw + gap)
        cy = y + sum(altos[:i // cols]) if vertical or cols == 1 else y
        if not (vertical or cols == 1):
            cy = y + (i // cols) * (alto_fila + gap)
        altos.append(pintar_tarjeta(slide, h, cx, cy, cw, ctx))
    if vertical or cols == 1:
        return sum(altos, Emu(0)) + Inches(0.1)
    return Emu(int(max(altos) if filas == 1 else (alto_fila + gap) * filas)) + Inches(0.1)


def pintar_tarjeta(slide, nodo, x, y, w, ctx):
    """.num / .paso / .panel de una grilla: marco + contenido."""
    if nodo.tiene('num'):
        etiqueta = nodo.buscar('span')
        valor = nodo.buscar('b')
        nota = nodo.buscar('i')
        alto = Inches(1.0) + (Inches(0.28) if nota else Emu(0))
        caja_texto(slide, x, y, w, alto, C['card'], C['line'])
        cy = y + Inches(0.14)
        if etiqueta:
            texto(slide, x + Inches(0.14), cy, w - Inches(0.28), Inches(0.2),
                  [(etiqueta[0].plano().upper(), True, C['muted'], True)], tam=8)
            cy += Inches(0.24)
        if valor:
            col = C['acc'] if valor[0].tiene('ok') else (
                C['mal'] if valor[0].tiene('mal') else C['ink'])
            texto(slide, x + Inches(0.14), cy, w - Inches(0.28), Inches(0.42),
                  [(valor[0].plano(), True, col, False)], tam=22)
            cy += Inches(0.46)
        if nota:
            texto(slide, x + Inches(0.14), cy, w - Inches(0.28), Inches(0.36),
                  [(nota[0].plano(), False, C['muted'], False)], tam=9.5)
        return alto + Inches(0.1)
    if nodo.tiene('paso'):
        alto = Inches(1.0)
        acc = nodo.tiene('acc')
        caja_texto(slide, x, y, w, alto, C['verdeclaro'] if acc else C['card'],
                   C['acc'] if acc else C['line'])
        texto(slide, x + Inches(0.1), y + Inches(0.12), w - Inches(0.2), alto - Inches(0.24),
              nodo.runs(), tam=11, alin=PP_ALIGN.CENTER, ancla=MSO_ANCHOR.MIDDLE)
        return alto + Inches(0.1)
    return pintar_contenedor(slide, nodo, x, y, w, ctx, marco=nodo.tiene('panel'),
                             alerta=nodo.tiene('alerta'))


# ------------------------------------------------------------------ diapositivas

def fondo(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def diapo_portada(prs, sec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fondo(s, C['bg'])
    rect(s, MARGEN, Inches(1.9), Inches(1.15), Inches(0.09), C['acc'], radio=True)
    h1 = sec.buscar('h1')[0]
    em = h1.buscar('em')
    titulo = [(t, n, c, m) for t, n, c, m in h1.runs()
              if not (em and t.strip() and t.strip() in em[0].plano())]
    tit = re.sub(r'\s+', ' ', ''.join(t for t, *_ in titulo)).strip()
    ht = alto_texto([(tit, True, None, False)], UTIL / 914400, 40, 1.06)
    texto(s, MARGEN, Inches(2.2), UTIL, ht, [(tit, True, C['ink'], False)], tam=40, espacio=1.06)
    if em:
        texto(s, MARGEN, Inches(2.2) + ht + Inches(0.22), UTIL, Inches(0.5),
              [(em[0].plano(), False, C['muted'], False)], tam=20)
    eq = sec.buscar('p', 'equipo')
    if eq:
        rect(s, MARGEN, Inches(6.2), UTIL, Inches(0.012), C['line'])
        nombres = eq[0].plano()
        sub = eq[0].buscar('em')
        if sub:
            nombres = nombres.replace(sub[0].plano(), '').strip()
        texto(s, MARGEN, Inches(6.4), UTIL, Inches(0.3),
              [(nombres, True, C['muted'], True)], tam=12)
        if sub:
            texto(s, MARGEN, Inches(6.75), UTIL, Inches(0.3),
                  [(sub[0].plano(), False, RGBColor(0x9A, 0xA8, 0xB4), True)], tam=10)
    return s


def diapo_seccion(prs, sec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fondo(s, C['dark'])
    n = sec.buscar('p', 'secn')
    h2 = sec.buscar('h2')[0]
    spec = sec.buscar('dl')
    y = Inches(1.15) if spec else Inches(2.0)
    if n:
        texto(s, MARGEN, y, UTIL, Inches(1.3), [(n[0].plano(), True, C['darkline'], True)],
              tam=76, espacio=0.85)
        y += Inches(1.25)
    texto(s, MARGEN, y, UTIL, Inches(0.95),
          [(h2.plano(), True, RGBColor(0xF4, 0xF7, 0xF6), False)], tam=40)
    y += Inches(0.95)
    rect(s, MARGEN, y, Inches(1.15), Inches(0.09), C['acc'], radio=True)
    y += Inches(0.32)
    for p in sec.buscar('p'):
        if p.tiene('secn'):
            continue
        if p.tiene('secbase'):
            y += Inches(0.16)
            texto(s, MARGEN, y, UTIL, Inches(0.8), p.runs(), tam=11, color=C['darkmuted'])
            y += Inches(0.5)
        else:
            texto(s, MARGEN, y, UTIL, Inches(0.42), p.runs(), tam=14, color=C['darkmuted'])
            y += Inches(0.42)
    if spec:
        pintar_spec(s, spec[0], MARGEN, y + Inches(0.1), UTIL)
    return s


def diapo_conclusiones(prs, sec, eyebrow):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fondo(s, C['bg'])
    texto(s, MARGEN, TOPE, UTIL, Inches(0.24), [(eyebrow, True, C['acc'], True)], tam=9)
    items = [li for li in sec.buscar('li')]
    y = Inches(0.95)
    alto = (ALTO - y - PIE) // max(1, len(items))
    for li in items:
        num = li.buscar('span', 'cn')
        h3 = li.buscar('h3')
        p = li.buscar('p')
        rect(s, MARGEN, y, UTIL, Inches(0.012), C['line'])
        if num:
            texto(s, MARGEN, y + Inches(0.12), Inches(0.75), Inches(0.5),
                  [(num[0].plano(), True, C['acc'], True)], tam=24)
        if h3:
            texto(s, MARGEN + Inches(0.85), y + Inches(0.1), UTIL - Inches(0.9), Inches(0.3),
                  [(h3[0].plano(), True, C['ink'], False)], tam=15)
        if p:
            texto(s, MARGEN + Inches(0.85), y + Inches(0.42), UTIL - Inches(0.9),
                  alto - Inches(0.5), p[0].runs(), tam=11.5)
        y += alto
    return s


def diapo_normal(prs, sec, eyebrow):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fondo(s, C['bg'])
    texto(s, MARGEN, TOPE, UTIL, Inches(0.24), [(eyebrow, True, C['acc'], True)], tam=9)
    y = Inches(0.82)
    h2 = sec.buscar('h2')
    if h2:
        runs = [(h2[0].plano(), True, C['ink'], False)]
        hh = alto_texto(runs, UTIL / 914400, 26)
        texto(s, MARGEN, y, UTIL, hh, runs, tam=26, espacio=1.1)
        y += hh + Inches(0.22)

    pie = sec.buscar('p', 'pieexp')
    tope_pie = ALTO - PIE - (Inches(0.55) if pie else Emu(0))
    for d in sec.buscar('p', 'destacado'):  # va debajo de la figura: reservarle el alto
        tope_pie -= alto_texto(d.runs(), UTIL / 914400 - 0.5, 12.5) + Inches(0.4)
    ctx = {'tmp': None, 'n': 0, 'css': '', 'max_fig': tope_pie - y - Inches(0.1)}
    ctx.update(prs._ctx)
    ctx['max_fig'] = tope_pie - y - Inches(0.1)

    cuerpo = [h for h in sec.hijos if isinstance(h, Nodo)
              and h.tag not in ('h2', 'h1') and not h.tiene('eyebrow')
              and not h.tiene('pieexp')]
    y0, marca = y, len(s.shapes)
    for nodo in cuerpo:
        if nodo.tiene('cols2') or nodo.tiene('interp'):
            cols = [h for h in nodo.hijos if isinstance(h, Nodo)]
            gap = Inches(0.34)
            if nodo.tiene('figura') and len(cols) == 2:
                anchos = [(UTIL - gap) * 62 // 100, (UTIL - gap) * 38 // 100]
            elif nodo.tiene('interp') and len(cols) == 2:
                anchos = [(UTIL - gap) * 33 // 100, (UTIL - gap) * 67 // 100]
            else:
                anchos = [(UTIL - gap) // len(cols)] * len(cols)
            cx, mayor = MARGEN, Emu(0)
            for col, aw in zip(cols, anchos):
                ctx['max_fig'] = tope_pie - y - Inches(0.1)
                mayor = max(mayor, pintar(s, col, cx, y, aw, ctx))
                cx += aw + gap
            y += mayor
        else:
            y += pintar(s, nodo, MARGEN, y, UTIL, ctx)

    sobra = (tope_pie - y0) - (y - y0)      # el HTML centra el cuerpo en vertical
    if sobra > Inches(0.25):
        for sh in list(s.shapes)[marca:]:
            if sh.top is not None:
                sh.top = sh.top + sobra // 2

    if pie:
        yb = ALTO - PIE - Inches(0.42)
        rect(s, MARGEN, yb, UTIL, Inches(0.01), C['line'])
        base = pie[0].buscar('span', 'pe-base')
        var = pie[0].buscar('span', 'pe-var')
        fin = pie[0].buscar('span', 'pe-fin')   # el chip de cierre va en verde, no en lila
        var = var or fin
        t = var[0].plano() if var else ''
        cw = Inches(0.075 * len(t) + 0.24) if var else Emu(0)
        texto(s, MARGEN, yb + Inches(0.1), Inches(0.55), Inches(0.2),
              [('BASE', True, C['muted'], True)], tam=8)
        if base:  # el ancho tiene que frenar antes de la pildora, si no se solapan
            texto(s, MARGEN + Inches(0.62), yb + Inches(0.08),
                  UTIL - cw - Inches(0.85), Inches(0.36),
                  [(base[0].plano(), True, RGBColor(0x93, 0xA2, 0xAD), True)], tam=8)
        if var:
            relleno = C['verdeclaro'] if fin else C['lila']
            borde = C['acc'] if fin else RGBColor(0xD6, 0xCD, 0xF0)
            tinta = RGBColor(0x0B, 0x7A, 0x63) if fin else C['acc2']
            caja_texto(s, MARGEN + UTIL - cw, yb + Inches(0.06), cw, Inches(0.28),
                       relleno, borde)
            texto(s, MARGEN + UTIL - cw, yb + Inches(0.06), cw, Inches(0.28),
                  [(t, True, tinta, True)], tam=8, alin=PP_ALIGN.CENTER,
                  ancla=MSO_ANCHOR.MIDDLE)
    return s


# ------------------------------------------------------------------ notas y main

def limpiar(md):
    md = re.sub(r'\*\*(.+?)\*\*', r'\1', md, flags=re.S)
    md = md.replace('`', '')
    md = re.sub(r'^---\s*$', '', md, flags=re.M)
    md = re.sub(r'^#+\s*', '', md, flags=re.M)
    md = re.sub(r'^> ?', '', md, flags=re.M)
    md = re.sub(r'^\|.*\|\s*$', '', md, flags=re.M)
    return re.sub(r'\n{3,}', '\n\n', md).strip()


def notas(n):
    partes = re.split(r'^(?=## \d+ · )', GUION.read_text(encoding='utf-8'), flags=re.M)
    intro, out = partes[0], {}
    for sec in partes[1:]:
        k = int(re.match(r'## (\d+)', sec).group(1))
        cuerpo, _, ap = sec.partition('\n# ')
        out[k] = limpiar(cuerpo)
        if ap:
            out[n] = out.get(n, '') + '\n\n' + limpiar('# ' + ap)
    out[1] = limpiar(intro) + '\n\n' + out.get(1, '')
    return out


def main():
    global SANS, MONO
    ap = argparse.ArgumentParser(description=__doc__.split('\n')[0])
    ap.add_argument('--fuentes', choices=sorted(FUENTES), default='office',
                    help='office: Arial/Consolas (default) · mac: Helvetica/Menlo, '
                         'identico al HTML en macOS · portable: Arial/Courier New')
    SANS, MONO = FUENTES[ap.parse_args().fuentes]

    prs = Presentation()
    prs.slide_width, prs.slide_height = ANCHO, ALTO
    css = re.search(r'<style>(.*?)</style>', HTML.read_text(encoding='utf-8'), re.S).group(1)
    secs = leer_slides()
    textos = notas(len(secs))
    with tempfile.TemporaryDirectory() as tmp:
        prs._ctx = {'tmp': Path(tmp), 'n': 0, 'css': css}
        for i, sec in enumerate(secs, 1):
            eb = sec.buscar('div', 'eyebrow')
            eyebrow = eb[0].plano() if eb else ''
            if sec.buscar('h1'):
                s = diapo_portada(prs, sec)
                eyebrow = 'portada'
            elif sec.tiene('sec'):
                s = diapo_seccion(prs, sec)
                eyebrow = sec.buscar('h2')[0].plano()
            elif sec.buscar('ul', 'conclu'):
                s = diapo_conclusiones(prs, sec, eyebrow)
            else:
                s = diapo_normal(prs, sec, eyebrow)
            s.notes_slide.notes_text_frame.text = textos.get(i, '')
            print(f'  {i:02d}  {eyebrow}')
        cp = prs.core_properties
        cp.title = 'Predicción de Buy Through Rate con Transformers'
        cp.subject = 'TP1 · 73.69 Large Language Models'
        prs.save(SALIDA)
    print(f'escrito {SALIDA.name} ({SALIDA.stat().st_size / 2**20:.1f} MB, {len(secs)} '
          f'diapositivas editables)')


if __name__ == '__main__':
    main()
