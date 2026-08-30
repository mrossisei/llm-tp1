"""Exporta presentacion.html a presentacion.pptx SIN alterar el diseno.

Cada diapositiva se renderiza con Chrome headless al tamano de proyeccion (1920x1080 por la
escala: 2 => 3840x2160) y entra al .pptx como imagen a pantalla completa (16:9); el guion.md va
como notas del orador (vista de presentador). El HTML sigue siendo la fuente (generar.py); este
script solo produce la version .pptx a partir de el. Requiere python-pptx y google-chrome.

    .venv/bin/python entrega/presentacion/exportar_pptx.py [--escala 2] [--png-dir carpeta]
"""
import argparse
import re
import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

AQUI = Path(__file__).resolve().parent
HTML = AQUI / 'presentacion.html'
GUION = AQUI / 'guion.md'
SALIDA = AQUI / 'presentacion.pptx'


MAC_CHROME = ('/Applications/Google Chrome.app/Contents/MacOS/Google Chrome',
              '/Applications/Chromium.app/Contents/MacOS/Chromium')


def chrome():
    for b in ('google-chrome', 'google-chrome-stable', 'chromium', 'chromium-browser'):
        if shutil.which(b):
            return b
    for ruta in MAC_CHROME:  # macOS no deja el binario en el PATH
        if Path(ruta).exists():
            return ruta
    raise SystemExit('no encuentro google-chrome ni chromium')


def renderizar(binario, n, destino, escala):
    """Captura la diapositiva n (la pagina muestra #sN) como PNG de 1920x1080 * escala."""
    subprocess.run([binario, '--headless=new', '--disable-gpu', '--hide-scrollbars',
                    '--window-size=1920,1080', f'--force-device-scale-factor={escala}',
                    '--virtual-time-budget=5000', f'--screenshot={destino}',
                    f'file://{HTML}#s{n}'], check=True, capture_output=True)
    if not destino.exists():
        raise SystemExit(f'Chrome no escribio {destino}')


def limpiar(md):
    """Markdown del guion -> texto plano legible en la vista de presentador."""
    md = re.sub(r'\*\*(.+?)\*\*', r'\1', md, flags=re.S)
    md = md.replace('`', '')
    md = re.sub(r'^---\s*$', '', md, flags=re.M)
    md = re.sub(r'^#+\s*', '', md, flags=re.M)
    md = re.sub(r'^> ?', '', md, flags=re.M)          # citas del guion
    md = re.sub(r'^\|.*\|\s*$', '', md, flags=re.M)  # tablas: ilegibles en la vista de notas
    return re.sub(r'\n{3,}', '\n\n', md).strip()


def rotular(html):
    """Un rotulo por diapositiva, en orden. Usa el eyebrow; si no hay (portada, divisorias),
    cae al <h1>/<h2>. Sirve de texto alternativo de la imagen en el .pptx."""
    etiquetas = []
    for sec in re.findall(r'<section class="slide[^"]*">(.*?)</section>', html, re.S):
        m = (re.search(r'<div class="eyebrow">([^<]*)</div>', sec)
             or re.search(r'<h2[^>]*>(.*?)</h2>', sec, re.S)
             or re.search(r'<h1[^>]*>(.*?)</h1>', sec, re.S))
        etiquetas.append(re.sub(r'\s+', ' ', re.sub(r'<[^>]+>', ' ', m.group(1))).strip()
                         if m else 'portada')
    return etiquetas


def notas(n):
    """Seccion '## k · ...' del guion -> notas de la diapositiva k. El apendice de preguntas va
    a la diapositiva de backup (la ultima) y la introduccion (duracion, recortes) a la portada."""
    partes = re.split(r'^(?=## \d+ · )', GUION.read_text(), flags=re.M)
    intro, secciones = partes[0], partes[1:]
    out = {}
    for sec in secciones:
        k = int(re.match(r'## (\d+)', sec).group(1))
        cuerpo, _, apendice = sec.partition('\n### ')
        out[k] = limpiar(cuerpo)
        if apendice:
            texto = limpiar('### ' + apendice)
            if n > k:
                out[n] = texto
            else:
                out[k] += '\n\n' + texto
    out[1] = limpiar(intro) + '\n\n' + out.get(1, '')
    return out


def main():
    ap = argparse.ArgumentParser(description=__doc__.split('\n')[0])
    ap.add_argument('--escala', type=int, default=2, help='factor de resolucion sobre 1920x1080')
    ap.add_argument('--png-dir', type=Path, help='conservar los PNG renderizados en esta carpeta')
    args = ap.parse_args()

    html = HTML.read_text()
    equipo = re.search(r'<p class="equipo">([^<]*)</p>', html)
    etiquetas = rotular(html)  # una por diapositiva, incluidas portada y divisorias
    n = len(etiquetas)
    binario = chrome()
    textos = notas(n)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)  # 16:9
    blanco = prs.slide_layouts[6]
    with tempfile.TemporaryDirectory() as tmp:
        carpeta = args.png_dir or Path(tmp)
        carpeta.mkdir(parents=True, exist_ok=True)
        for i, etiqueta in enumerate(etiquetas, 1):
            png = carpeta / f's{i:02d}.png'
            renderizar(binario, i, png, args.escala)
            s = prs.slides.add_slide(blanco)
            pic = s.shapes.add_picture(str(png), 0, 0, prs.slide_width, prs.slide_height)
            pic.name = f'diapo {i:02d}'
            pic._element._nvXxPr.cNvPr.set('descr', etiqueta)  # texto alternativo
            s.notes_slide.notes_text_frame.text = textos.get(i, '')
            print(f'  {i:02d}  {etiqueta}')
        cp = prs.core_properties
        cp.title = 'Predicción de Buy Through Rate con Transformers'
        cp.subject = 'TP1 · 73.69 Large Language Models'
        cp.author = cp.last_modified_by = equipo.group(1) if equipo else ''
        prs.save(SALIDA)
    print(f'escrito {SALIDA.name} ({SALIDA.stat().st_size / 2**20:.1f} MB, {n} diapositivas, '
          f'notas desde {GUION.name})')


if __name__ == '__main__':
    main()
